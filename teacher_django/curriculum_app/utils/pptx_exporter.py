import os
import json
import time
from django.conf import settings


def export_outline_to_pptx(outline, progress_callback=None):
    """
    将 CourseOutline 导出为 PPTX 文件，返回 (file_path, filename)

    实现要点：
    - 使用 python-pptx 生成幻灯片，尽量增加视觉效果（背景图、标题样式、要点分步展示模拟动画）。
    - 若可用，会尝试从 Unsplash 拉取主题相关图片作为幻灯片背景，失败则使用纯色背景。
    - 最终文件保存在 MEDIA_ROOT/exports/ 下，并返回文件路径与文件名。
    """
    # 延迟导入依赖，便于在缺少依赖时捕获错误
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
        from pptx.enum.text import PP_ALIGN
    except Exception as e:
        raise RuntimeError('缺少 python-pptx 库，请安装：pip install python-pptx Pillow')

    try:
        od = json.loads(outline.outline_data)
    except Exception:
        od = {}

    title = od.get('title') or outline.title or '课程'

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def paint_solid(shape, color, transparency=0):
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = color
        fill.transparency = transparency
        shape.line.fill.background()

    def add_shape_box(target_slide, left, top, width, height, fill, radius=True, transparency=0, line=None):
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
        shape = target_slide.shapes.add_shape(shape_type, left, top, width, height)
        paint_solid(shape, fill, transparency=transparency)
        if line:
            shape.line.color.rgb = line
            shape.line.width = Pt(1)
        return shape

    def add_rich_text(target_slide, text, left, top, width, height, size=18, bold=False, color=None, align=None, valign=None):
        box = target_slide.shapes.add_textbox(left, top, width, height)
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.margin_left = Inches(0.05)
        frame.margin_right = Inches(0.05)
        frame.margin_top = Inches(0.02)
        frame.margin_bottom = Inches(0.02)
        if valign:
            frame.vertical_anchor = valign
        paragraph = frame.paragraphs[0]
        paragraph.text = str(text or '')
        paragraph.font.name = 'Microsoft YaHei'
        paragraph.font.size = Pt(size)
        paragraph.font.bold = bold
        if color:
            paragraph.font.color.rgb = color
        if align:
            paragraph.alignment = align
        return box

    def add_deck_cover():
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = RGBColor(255, 255, 255)
        ink = RGBColor(18, 24, 38)
        muted = RGBColor(85, 95, 112)
        accent = RGBColor(18, 86, 138)
        accent2 = RGBColor(230, 148, 46)
        deep = RGBColor(245, 245, 245)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = bg
        add_shape_box(slide, Inches(0), Inches(0), Inches(13.333), Inches(7.5), bg, radius=False)
        add_shape_box(slide, Inches(0), Inches(0), Inches(4.35), Inches(7.5), deep, radius=False)
        add_shape_box(slide, Inches(4.35), Inches(0), Inches(0.08), Inches(7.5), accent2, radius=False)
        add_shape_box(slide, Inches(10.55), Inches(0.55), Inches(1.95), Inches(1.95), accent, transparency=0)
        add_shape_box(slide, Inches(11.2), Inches(4.62), Inches(1.05), Inches(1.05), accent2, transparency=0)
        add_shape_box(slide, Inches(9.48), Inches(5.42), Inches(0.55), Inches(0.55), RGBColor(36, 161, 148), transparency=0)
        add_rich_text(slide, 'AI COURSEWARE', Inches(0.82), Inches(0.78), Inches(2.65), Inches(0.28), size=11, bold=True, color=ink)
        add_rich_text(slide, '讲解课件', Inches(0.82), Inches(1.16), Inches(2.15), Inches(0.36), size=16, bold=True, color=accent2)
        add_rich_text(slide, title, Inches(4.98), Inches(1.42), Inches(6.45), Inches(1.75), size=42, bold=True, color=ink)
        subtitle = f"由智能助手生成 · {getattr(outline.user, 'username', 'teacher')}"
        add_rich_text(slide, subtitle, Inches(5.03), Inches(3.44), Inches(5.4), Inches(0.36), size=14, color=muted)
        for idx, label in enumerate(['结构化讲解', 'H5 演示', '随堂测验']):
            left = Inches(5.02 + idx * 1.75)
            add_shape_box(slide, left, Inches(4.55), Inches(1.28), Inches(0.42), accent if idx == 0 else (RGBColor(36, 161, 148) if idx == 1 else accent2), transparency=0)
            add_rich_text(slide, label, left, Inches(4.66), Inches(1.28), Inches(0.12), size=8, bold=True, color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER)
        add_shape_box(slide, Inches(0.82), Inches(5.35), Inches(2.8), Inches(0.02), accent2, radius=False)
        add_rich_text(slide, 'GENERATED DECK', Inches(0.82), Inches(5.62), Inches(2.6), Inches(0.22), size=10, color=RGBColor(150, 150, 150))
        add_rich_text(slide, '面向课堂讲解与自学复盘', Inches(0.82), Inches(6.02), Inches(2.8), Inches(0.28), size=11, color=muted)
        return slide

    add_deck_cover()

    # Parse ppt preview into slides
    ppt_preview = (od.get('resources', {}) or {}).get('ppt', {}) or {}
    structured = ppt_preview.get('structured_slides') or od.get('slide_deck') or ppt_preview.get('slides') or od.get('slides')
    structured_slides = structured if isinstance(structured, list) and all(isinstance(item, dict) for item in structured) else []
    if not structured_slides:
        from agent_system.generation import normalize_slide_deck
        structured_slides = normalize_slide_deck(None, title, od)

    if structured_slides:
        palettes = {
            'academic_light': {
                'bg': RGBColor(255, 255, 255), 'panel': RGBColor(255, 255, 255), 'ink': RGBColor(18, 24, 38),
                'muted': RGBColor(83, 96, 117), 'accent': RGBColor(18, 86, 138), 'accent2': RGBColor(230, 148, 46),
                'line': RGBColor(218, 224, 232), 'soft': RGBColor(235, 242, 248), 'shadow': RGBColor(228, 232, 238),
                'rail': RGBColor(245, 245, 245), 'rail_text': RGBColor(18, 24, 38), 'success': RGBColor(36, 161, 148),
            },
            'tech_blue': {
                'bg': RGBColor(255, 255, 255), 'panel': RGBColor(255, 255, 255), 'ink': RGBColor(15, 28, 44),
                'muted': RGBColor(61, 77, 96), 'accent': RGBColor(11, 105, 135), 'accent2': RGBColor(59, 130, 246),
                'line': RGBColor(196, 224, 239), 'soft': RGBColor(224, 244, 248), 'shadow': RGBColor(218, 231, 239),
                'rail': RGBColor(245, 245, 245), 'rail_text': RGBColor(15, 28, 44), 'success': RGBColor(14, 165, 156),
            },
            'chalkboard_dark': {
                'bg': RGBColor(255, 255, 255), 'panel': RGBColor(255, 255, 255), 'ink': RGBColor(18, 24, 38),
                'muted': RGBColor(83, 96, 117), 'accent': RGBColor(67, 214, 159), 'accent2': RGBColor(245, 184, 61),
                'line': RGBColor(218, 224, 232), 'soft': RGBColor(235, 242, 248), 'shadow': RGBColor(228, 232, 238),
                'rail': RGBColor(245, 245, 245), 'rail_text': RGBColor(18, 24, 38), 'success': RGBColor(36, 161, 148),
            },
        }

        def set_background(target_slide, palette):
            fill = target_slide.background.fill
            fill.solid()
            fill.fore_color.rgb = palette['bg']
            add_shape_box(target_slide, Inches(0), Inches(0), Inches(1.02), Inches(7.5), palette['rail'], radius=False)
            add_shape_box(target_slide, Inches(1.02), Inches(0), Inches(0.05), Inches(7.5), palette['accent2'], radius=False)

        def add_textbox(target_slide, text, left, top, width, height, size=18, bold=False, color=None, align=None):
            return add_rich_text(target_slide, text, left, top, width, height, size=size, bold=bold, color=color, align=align)

        def add_card(target_slide, title_text, items, left, top, width, height, palette, index=None, dense=False):
            add_shape_box(target_slide, left + Inches(0.04), top + Inches(0.05), width, height, palette['shadow'], transparency=42)
            card = add_shape_box(target_slide, left, top, width, height, palette['panel'], transparency=0, line=palette['line'])
            if index is not None:
                add_shape_box(target_slide, left, top, Inches(0.16), height, palette['accent'], radius=False)
                add_shape_box(target_slide, left + Inches(0.34), top + Inches(0.25), Inches(0.5), Inches(0.5), palette['soft'], transparency=0)
                add_textbox(target_slide, str(index), left + Inches(0.34), top + Inches(0.41), Inches(0.5), Inches(0.14), size=8, bold=True, color=palette['accent'], align=PP_ALIGN.CENTER)
                title_left = left + Inches(1.0)
                title_width = width - Inches(1.22)
            else:
                add_shape_box(target_slide, left, top, width, Inches(0.12), palette['accent'], radius=False)
                title_left = left + Inches(0.36)
                title_width = width - Inches(0.7)
            add_textbox(target_slide, title_text or '要点', title_left, top + Inches(0.24), title_width, Inches(0.42), size=15 if not dense else 13, bold=True, color=palette['ink'])
            return card

        def add_bullet_box(target_slide, title_text, items, left, top, width, height, palette, index=None, dense=False):
            add_card(target_slide, title_text, items, left, top, width, height, palette, index=index, dense=dense)
            box = target_slide.shapes.add_textbox(left + Inches(0.36), top + Inches(0.86), width - Inches(0.66), height - Inches(1.05))
            frame = box.text_frame
            frame.clear()
            frame.word_wrap = True
            frame.margin_left = 0
            frame.margin_right = 0
            frame.margin_top = 0
            frame.margin_bottom = 0
            clean_items = [str(item or '').strip() for item in items if str(item or '').strip()]
            for item_index, item in enumerate(clean_items[:5]):
                para = frame.paragraphs[0] if item_index == 0 else frame.add_paragraph()
                para.text = '• ' + str(item or '')
                para.level = 0
                para.font.name = 'Microsoft YaHei'
                para.font.size = Pt(12 if dense else 14)
                para.font.color.rgb = palette['muted']
                para.space_after = Pt(5)
            return box

        def block_items(block):
            items = block.get('items') or block.get('points') or []
            if not isinstance(items, list):
                items = [items]
            text = block.get('text') or block.get('content')
            if text:
                items = [text] + items
            return [str(item) for item in items if str(item or '').strip()]

        def add_code_box(target_slide, block, left, top, width, height):
            bg_color = RGBColor(40, 44, 52)
            text_color = RGBColor(171, 178, 191)
            header_color = RGBColor(33, 37, 43)
            label = str(block.get('label') or '代码示例').strip()
            language = str(block.get('language') or '').strip()

            header_height = Inches(0.4)
            add_shape_box(target_slide, left, top, width, header_height, header_color, radius=False)
            header_text = f'{label}（{language}）' if language else label
            add_textbox(target_slide, header_text, left + Inches(0.18), top + Inches(0.06), width - Inches(0.3), Inches(0.28), size=11, bold=True, color=RGBColor(229, 192, 123))

            code_top = top + header_height
            code_height = height - header_height
            add_shape_box(target_slide, left, code_top, width, code_height, bg_color, radius=False)

            code = str(block.get('code') or '').rstrip('\n')
            lines = code.splitlines() or ['（暂无代码内容）']
            max_lines = 18
            if len(lines) > max_lines:
                lines = lines[:16] + ['...完整代码见网页课件']

            box = target_slide.shapes.add_textbox(left + Inches(0.16), code_top + Inches(0.08), width - Inches(0.3), code_height - Inches(0.16))
            frame = box.text_frame
            frame.clear()
            frame.word_wrap = False
            frame.margin_left = 0
            frame.margin_right = 0
            frame.margin_top = 0
            frame.margin_bottom = 0
            font_size = 11 if len(lines) <= 12 else 9
            max_chars = 90
            for line_index, line in enumerate(lines):
                para = frame.paragraphs[0] if line_index == 0 else frame.add_paragraph()
                display = (line[:max_chars] + '..') if len(line) > max_chars else line
                para.text = display if display.strip() else ' '
                para.font.name = 'Consolas'
                para.font.size = Pt(font_size)
                para.font.color.rgb = text_color
                para.space_after = Pt(0)

        def add_quiz_box(target_slide, blocks, left, top, width, height, palette):
            # Accept both question_text and text fields; handle MC and open-ended
            questions = [
                b for b in (blocks or [])
                if isinstance(b, dict) and b.get('kind') == 'question'
                and (b.get('question_text') or b.get('text'))
            ]
            if not questions:
                items = []
                for b in (blocks or []):
                    items.extend(block_items(b))
                add_bullet_box(target_slide, '随堂检查', items or ['请用一句话解释本页核心概念。'],
                               left, top, width, height, palette)
                return

            if any(q.get('choices') for q in questions):
                # --- Multiple-choice question ---
                q = next(q for q in questions if q.get('choices'))
                q_text = str(q.get('question_text') or q.get('text') or '').strip()
                raw_correct = str(q.get('correct_answer') or '').strip()
                choices = q.get('choices') or []

                # Parse each choice into (letter_value, display_label)
                parsed = []
                for c in choices:
                    if isinstance(c, dict):
                        c_lbl = str(c.get('label') or c.get('value') or '').strip()
                        c_val = str(c.get('value') or '').strip().upper()
                        if len(c_val) > 2:
                            c_val = c_lbl[:1].upper() if c_lbl else ''
                    else:
                        c_lbl = str(c).strip()
                        c_val = c_lbl[:1].upper() if c_lbl else ''
                    parsed.append((c_val, c_lbl))

                # Resolve correct_answer: handle both letter and full-text formats
                correct = raw_correct.strip().upper()
                if len(correct) > 2:
                    resolved = False
                    for c_val, c_lbl in parsed:
                        if c_val and (correct.startswith(c_val) or raw_correct.strip() in c_lbl):
                            correct, resolved = c_val, True
                            break
                    if not resolved:
                        correct = correct[:1]

                ct = top + Inches(0.14)
                add_textbox(target_slide, f"Q：{q_text}", left + Inches(0.20), ct,
                            width - Inches(0.40), Inches(0.56), size=15, bold=True, color=palette['ink'])
                ct += Inches(0.62)

                col_w = (width - Inches(0.52)) / 2
                l_col = left + Inches(0.20)
                r_col = l_col + col_w + Inches(0.12)
                ch = Inches(0.44)
                rows = (len(parsed) + 1) // 2
                for row in range(rows):
                    for col in range(2):
                        ci = row * 2 + col
                        if ci >= len(parsed):
                            break
                        c_val, c_lbl = parsed[ci]
                        ok = bool(c_val) and c_val == correct
                        cx = l_col if col == 0 else r_col
                        cy = ct + row * (ch + Inches(0.06))
                        bg = palette['soft'] if ok else palette['shadow']
                        add_shape_box(target_slide, cx - Inches(0.04), cy - Inches(0.02),
                                      col_w + Inches(0.04), ch, bg, transparency=25)
                        add_textbox(target_slide, ('✓ ' if ok else '') + c_lbl, cx, cy,
                                    col_w - Inches(0.04), ch, size=12, bold=ok,
                                    color=palette['success'] if ok else palette['muted'])
                ct += rows * (ch + Inches(0.06)) + Inches(0.10)

                exp = str(q.get('explanation') or '').strip()
                if exp:
                    exp = (exp[:148] + '...') if len(exp) > 150 else exp
                    exp_h = Inches(0.56)
                    add_shape_box(target_slide, left + Inches(0.20), ct,
                                  width - Inches(0.40), exp_h, palette['soft'], transparency=20)
                    add_textbox(target_slide, f"解析：{exp}", left + Inches(0.28), ct + Inches(0.04),
                                width - Inches(0.56), exp_h - Inches(0.08), size=11, color=palette['muted'])
            else:
                # --- Open-ended / short-answer questions ---
                qs = questions[:3]
                row_h = (height - Inches(0.28)) / max(len(qs), 1)
                badge_colors = [palette['accent'], palette['success'], palette['accent2']]
                for qi, q in enumerate(qs):
                    q_text = str(q.get('question_text') or q.get('text') or '').strip()
                    q_label = str(q.get('label') or f'思考 {qi + 1}').strip()
                    rt = top + Inches(0.14) + qi * row_h
                    add_shape_box(target_slide, left + Inches(0.20), rt + Inches(0.10),
                                  Inches(0.36), Inches(0.36), badge_colors[qi % 3])
                    add_textbox(target_slide, str(qi + 1), left + Inches(0.20),
                                rt + Inches(0.14), Inches(0.36), Inches(0.22),
                                size=11, bold=True, color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER)
                    add_textbox(target_slide, f'[{q_label}]', left + Inches(0.64),
                                rt + Inches(0.12), Inches(1.50), Inches(0.22),
                                size=9, bold=True, color=palette['accent'])
                    txt_h = max(Inches(0.30), row_h - Inches(0.46))
                    add_textbox(target_slide, q_text, left + Inches(0.64), rt + Inches(0.36),
                                width - Inches(0.84), txt_h, size=14,
                                bold=(len(qs) == 1), color=palette['ink'])

        def render_blocks(target_slide, slide_data, palette):
            layout = str(slide_data.get('layout') or slide_data.get('type') or 'two_column')
            blocks = slide_data.get('visual_blocks') or []
            if not isinstance(blocks, list) or not blocks:
                bullets = slide_data.get('bullets') or []
                if not isinstance(bullets, list):
                    bullets = [bullets]
                blocks = [{'label': f'要点 {idx}', 'text': item} for idx, item in enumerate(bullets[:4], start=1)]

            code_blocks = [b for b in blocks if isinstance(b, dict) and b.get('kind') in ('code', 'code_block') and b.get('code')]
            has_code = len(code_blocks) > 0

            if layout == 'cover':
                add_shape_box(target_slide, Inches(1.55), Inches(1.54), Inches(6.25), Inches(1.38), palette['panel'], line=palette['line'])
                add_shape_box(target_slide, Inches(1.55), Inches(1.54), Inches(0.18), Inches(1.38), palette['accent2'], radius=False)
                add_textbox(target_slide, slide_data.get('title') or title, Inches(1.96), Inches(1.88), Inches(5.45), Inches(0.46), size=23, bold=True, color=palette['ink'])
                add_bullet_box(target_slide, '本讲学习目标', slide_data.get('bullets') or [], Inches(1.58), Inches(3.35), Inches(5.68), Inches(2.56), palette)
                add_shape_box(target_slide, Inches(7.68), Inches(3.35), Inches(4.12), Inches(2.56), palette['rail'], transparency=0)
                add_textbox(target_slide, '课堂路径', Inches(8.02), Inches(3.76), Inches(1.8), Inches(0.34), size=15, bold=True, color=palette['rail_text'])
                for pos, label in enumerate(['讲解', '演示', '练习'], start=1):
                    left = Inches(8.02 + (pos - 1) * 1.1)
                    add_shape_box(target_slide, left, Inches(4.52), Inches(0.72), Inches(0.72), palette['accent2'] if pos == 2 else palette['success'])
                    add_textbox(target_slide, str(pos), left, Inches(4.72), Inches(0.72), Inches(0.18), size=10, bold=True, color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER)
                    add_textbox(target_slide, label, left - Inches(0.12), Inches(5.28), Inches(0.96), Inches(0.24), size=10, color=palette['rail_text'], align=PP_ALIGN.CENTER)
                return

            if layout in ('two_column', 'comparison'):
                columns = blocks[:2] if len(blocks) >= 2 else blocks + [{'label': '补充说明', 'items': slide_data.get('bullets') or []}]
                for idx, block in enumerate(columns[:2]):
                    col_left = Inches(1.58 + idx * 5.90)
                    col_top = Inches(2.02)
                    col_width = Inches(5.50)
                    col_height = Inches(3.78)
                    if block.get('kind') in ('code', 'code_block') and block.get('code'):
                        add_code_box(target_slide, block, col_left, col_top, col_width, col_height)
                    else:
                        add_bullet_box(target_slide, block.get('label') or f'栏 {idx + 1}', block_items(block), col_left, col_top, col_width, col_height, palette, index=idx + 1)
                if layout == 'comparison':
                    add_shape_box(target_slide, Inches(7.01), Inches(3.25), Inches(0.54), Inches(0.54), palette['accent2'])
                    add_textbox(target_slide, 'VS', Inches(7.02), Inches(3.42), Inches(0.52), Inches(0.12), size=8, bold=True, color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER)
                return

            if layout in ('process_flow', 'agenda', 'concept_map'):
                if layout == 'concept_map':
                    add_shape_box(target_slide, Inches(5.28), Inches(2.48), Inches(3.05), Inches(1.1), palette['accent'], transparency=0)
                    add_textbox(target_slide, slide_data.get('title') or '核心概念', Inches(5.5), Inches(2.84), Inches(2.62), Inches(0.24), size=15, bold=True, color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER)
                    positions = [(1.58, 1.72), (8.92, 1.72), (1.58, 4.28), (8.92, 4.28)]
                    for idx, block in enumerate(blocks[:4]):
                        left, top = positions[idx]
                        if block.get('kind') in ('code', 'code_block') and block.get('code'):
                            add_code_box(target_slide, block, Inches(left), Inches(top), Inches(3.45), Inches(1.72))
                        else:
                            add_bullet_box(target_slide, block.get('label') or f'概念 {idx + 1}', block_items(block), Inches(left), Inches(top), Inches(3.45), Inches(1.72), palette, dense=True)
                else:
                    usable = blocks[:5] or [{'label': '要点', 'items': slide_data.get('bullets') or []}]
                    width = 11.40 / max(len(usable), 1)
                    for idx, block in enumerate(usable):
                        left = Inches(1.56 + idx * width)
                        if block.get('kind') in ('code', 'code_block') and block.get('code'):
                            add_code_box(target_slide, block, left, Inches(2.42), Inches(max(width - 0.16, 1.64)), Inches(2.62))
                        else:
                            add_bullet_box(target_slide, block.get('label') or f'步骤 {idx + 1}', block_items(block), left, Inches(2.42), Inches(max(width - 0.16, 1.64)), Inches(2.62), palette, index=idx + 1, dense=len(usable) >= 4)
                        if idx < len(usable) - 1:
                            start_x = Inches(1.56 + idx * width + max(width - 0.16, 1.64) - 0.04)
                            line = target_slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, start_x, Inches(3.74), Inches(1.56 + (idx + 1) * width - 0.1), Inches(3.74))
                            line.line.color.rgb = palette['accent']
                            line.line.width = Pt(2)
                return

            if layout == 'case_study':
                main = blocks[0] if blocks else {'label': '案例', 'items': slide_data.get('bullets') or []}
                support = blocks[1:] or [{'label': '分析路径', 'items': slide_data.get('bullets') or []}]
                if main.get('kind') in ('code', 'code_block') and main.get('code'):
                    add_code_box(target_slide, main, Inches(1.58), Inches(2.02), Inches(11.37), Inches(3.78))
                else:
                    add_shape_box(target_slide, Inches(1.58), Inches(2.02), Inches(5.36), Inches(3.86), palette['rail'], transparency=0)
                    add_textbox(target_slide, main.get('label') or '案例情境', Inches(1.96), Inches(2.42), Inches(4.64), Inches(0.38), size=18, bold=True, color=palette['rail_text'])
                    case_items = block_items(main) or slide_data.get('bullets') or []
                    add_textbox(target_slide, case_items[0] if case_items else '用一个真实问题引出知识点。', Inches(1.96), Inches(3.05), Inches(4.58), Inches(1.46), size=20, bold=True, color=RGBColor(255, 255, 255))
                    add_textbox(target_slide, slide_data.get('student_interaction') or '请学生先做判断，再解释理由。', Inches(1.98), Inches(5.16), Inches(4.52), Inches(0.32), size=11, color=palette['rail_text'])
                    for idx, block in enumerate(support[:2]):
                        add_bullet_box(target_slide, block.get('label') or f'分析 {idx + 1}', block_items(block), Inches(7.32), Inches(2.02 + idx * 1.98), Inches(5.63), Inches(1.68), palette, index=idx + 1, dense=True)
                return

            if layout == 'summary':
                summary_items = slide_data.get('bullets') or []
                if not isinstance(summary_items, list):
                    summary_items = [summary_items]
                for idx, item in enumerate(summary_items[:3]):
                    left = Inches(1.58 + idx * 3.80)
                    add_shape_box(target_slide, left, Inches(2.24), Inches(3.56), Inches(2.72), palette['panel'], line=palette['line'])
                    add_shape_box(target_slide, left + Inches(0.28), Inches(2.52), Inches(0.58), Inches(0.58), palette['accent'] if idx == 0 else (palette['success'] if idx == 1 else palette['accent2']))
                    add_textbox(target_slide, str(idx + 1), left + Inches(0.28), Inches(2.69), Inches(0.58), Inches(0.14), size=8, bold=True, color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER)
                    add_textbox(target_slide, str(item), left + Inches(0.34), Inches(3.32), Inches(2.90), Inches(0.96), size=15, bold=True, color=palette['ink'], align=PP_ALIGN.CENTER)
                add_textbox(target_slide, slide_data.get('student_interaction') or '用一句话说出今天最重要的收获。', Inches(1.60), Inches(5.44), Inches(10.0), Inches(0.38), size=13, color=palette['muted'], align=PP_ALIGN.CENTER)
                return

            if layout == 'quiz_check':
                add_shape_box(target_slide, Inches(1.58), Inches(1.48), Inches(11.37), Inches(3.60), palette['panel'], line=palette['line'])
                has_mc = any(isinstance(b, dict) and b.get('kind') == 'question' and b.get('choices') for b in blocks)
                if has_mc:
                    add_shape_box(target_slide, Inches(1.98), Inches(1.86), Inches(0.96), Inches(0.96), palette['accent2'])
                    add_textbox(target_slide, '?', Inches(1.98), Inches(2.11), Inches(0.96), Inches(0.24), size=18, bold=True, color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER)
                    add_quiz_box(target_slide, blocks, Inches(3.28), Inches(1.64), Inches(9.37), Inches(3.28), palette)
                else:
                    add_quiz_box(target_slide, blocks, Inches(1.78), Inches(1.64), Inches(10.97), Inches(3.28), palette)
                return

            if layout == 'animation_embed':
                anim_block = next((b for b in blocks if isinstance(b, dict) and b.get('kind') == 'animation'), (blocks[0] if blocks else {}))
                concept_name = str(anim_block.get('label') or anim_block.get('concept_name') or slide_data.get('title') or '动态演示').strip()
                usage_note = str(anim_block.get('usage_note') or slide_data.get('visual_hint') or '观察动画的变化过程，理解其背后的原理。').strip()
                if len(usage_note) > 90:
                    usage_note = usage_note[:90] + '...'

                add_shape_box(target_slide, Inches(1.58), Inches(1.82), Inches(6.06), Inches(4.02), palette['panel'], line=palette['line'])
                add_shape_box(target_slide, Inches(1.9), Inches(2.18), Inches(5.42), Inches(2.98), palette['soft'], radius=False, transparency=0)
                add_shape_box(target_slide, Inches(4.36), Inches(3.0), Inches(0.6), Inches(0.6), palette['accent'])
                play = target_slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(4.55), Inches(3.16), Inches(0.26), Inches(0.26))
                paint_solid(play, RGBColor(255, 255, 255))
                play.rotation = 90
                add_textbox(target_slide, concept_name, Inches(2.1), Inches(3.78), Inches(5.0), Inches(0.36), size=14, bold=True, color=palette['ink'], align=PP_ALIGN.CENTER)
                add_textbox(target_slide, '请前往网页课件查看交互动画演示', Inches(2.1), Inches(4.18), Inches(5.0), Inches(0.5), size=11, color=palette['muted'], align=PP_ALIGN.CENTER)
                add_bullet_box(target_slide, '互动提示', [usage_note, slide_data.get('student_interaction') or '观察动画并回答关键问题。'], Inches(7.92), Inches(2.08), Inches(5.03), Inches(3.05), palette)
                return

            if has_code:
                add_code_box(target_slide, code_blocks[0], Inches(1.58), Inches(2.02), Inches(11.37), Inches(3.78))
            else:
                add_bullet_box(target_slide, '核心要点', slide_data.get('bullets') or [], Inches(1.58), Inches(2.08), Inches(11.37), Inches(3.62), palette)

        def add_slide_chrome(target_slide, slide_data, palette, index, total):
            layout = str(slide_data.get('layout') or slide_data.get('type') or 'slide')
            add_textbox(target_slide, f'{index:02d}', Inches(0.18), Inches(0.44), Inches(0.64), Inches(0.28), size=14, bold=True, color=palette['rail_text'], align=PP_ALIGN.CENTER)
            add_textbox(target_slide, layout.replace('_', ' ').upper(), Inches(1.56), Inches(0.42), Inches(2.8), Inches(0.22), size=8, bold=True, color=palette['accent'])
            if layout != 'cover':
                title_text = str(slide_data.get('title') or f'第 {index} 页').strip()
                if len(title_text) > 22:
                    title_size = 18
                elif len(title_text) > 16:
                    title_size = 21
                else:
                    title_size = 25
                if len(title_text) > 40:
                    title_text = title_text[:40] + '...'
                add_textbox(target_slide, title_text, Inches(1.56), Inches(0.7), Inches(9.8), Inches(0.62), size=title_size, bold=True, color=palette['ink'])
            add_textbox(target_slide, f'{index:02d}/{total:02d}', Inches(11.18), Inches(6.98), Inches(0.78), Inches(0.18), size=8, bold=True, color=palette['muted'], align=PP_ALIGN.RIGHT)
            add_textbox(target_slide, title, Inches(1.56), Inches(6.98), Inches(4.8), Inches(0.18), size=8, color=palette['muted'])

        total_steps = len(structured_slides) + 1
        for index, slide_data in enumerate(structured_slides, start=1):
            page = prs.slides.add_slide(prs.slide_layouts[6])
            palette = palettes.get(str(slide_data.get('theme') or 'academic_light'), palettes['academic_light'])
            set_background(page, palette)
            add_slide_chrome(page, slide_data, palette, index, len(structured_slides))
            render_blocks(page, slide_data, palette)
            notes = slide_data.get('speaker_notes')
            if notes:
                notes = str(notes).strip()
                if len(notes) > 120:
                    notes = notes[:120] + '...'
                add_shape_box(page, Inches(1.56), Inches(6.18), Inches(11.39), Inches(0.48), palette['soft'], transparency=12, line=palette['line'])
                add_textbox(page, f"讲稿：{notes}", Inches(1.78), Inches(6.26), Inches(10.89), Inches(0.36), size=8, color=palette['muted'])
            if progress_callback:
                try:
                    progress_callback(max(0, min(99, int((index / float(total_steps)) * 100))))
                except Exception:
                    pass

        exports_dir = os.path.join(settings.MEDIA_ROOT, 'exports')
        os.makedirs(exports_dir, exist_ok=True)
        filename = f'outline_{outline.id}_{int(time.time())}.pptx'
        file_path = os.path.join(exports_dir, filename)
        prs.save(file_path)
        if progress_callback:
            try:
                progress_callback(100)
            except Exception:
                pass
        return file_path, filename
