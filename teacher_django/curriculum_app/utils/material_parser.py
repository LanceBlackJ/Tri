import os
import re
from typing import Dict, List


class MaterialParseError(Exception):
    pass


MATERIAL_TYPE_EXT_MAP = {
    'pdf': 'pdf',
    'ppt': 'ppt', 'pptx': 'ppt',
    'doc': 'doc', 'docx': 'doc',
    'mp4': 'video', 'mov': 'video', 'avi': 'video', 'mkv': 'video', 'webm': 'video',
    'jpg': 'image', 'jpeg': 'image', 'png': 'image', 'gif': 'image', 'webp': 'image', 'bmp': 'image',
    'zip': 'archive', 'rar': 'archive', '7z': 'archive', 'tar': 'archive', 'gz': 'archive',
}


def infer_material_type_from_filename(filename):
    """按文件扩展名自动识别资料类型。上传表单与 backfill_material_types
    命令共用同一份规则，避免两处判断逻辑不一致。"""
    ext = os.path.splitext(filename or '')[1].lower().lstrip('.')
    return MATERIAL_TYPE_EXT_MAP.get(ext, 'other')


def title_from_filename(filename):
    """标题留空时，用文件名自动生成：去掉扩展名，下划线/短横线换成空格。"""
    stem = os.path.splitext(filename or '')[0]
    stem = re.sub(r'[_-]+', ' ', stem).strip()
    return stem


def _strip_invalid_surrogates(value):
    if not value:
        return ''
    return ''.join(ch for ch in str(value) if not 0xD800 <= ord(ch) <= 0xDFFF)


_WINGDINGS_BULLET_RE = re.compile(r'^p\s+(?=\S)')


def _fix_bullet_artifacts(text):
    """PPT/PDF导出文本中，Wingdings符号字体的项目符号常被提取为字面字符'p'，
    将形如'p 内容...'的行首替换为标准的'- '项目符号。"""
    lines = text.split('\n')
    fixed = []
    for line in lines:
        stripped = line.lstrip()
        leading_ws = line[:len(line) - len(stripped)]
        new_stripped = _WINGDINGS_BULLET_RE.sub('- ', stripped, count=1)
        fixed.append(leading_ws + new_stripped)
    return '\n'.join(fixed)


def _normalize_text(value):
    text = _strip_invalid_surrogates(value).replace('\x00', ' ').strip()
    return _fix_bullet_artifacts(text)


def _summarize_keywords(text: str, limit: int = 6) -> str:
    words = []
    seen = set()
    for raw in text.replace('\n', ' ').split(' '):
        token = raw.strip(' ,.;:!?' '"()[]{}')
        if len(token) < 2:
            continue
        lowered = token.lower()
        if lowered in seen:
            continue
        seen.add(lowered)
        words.append(token)
        if len(words) >= limit:
            break
    return ' / '.join(words)


def _build_chunk(index: int, source_page: str, heading: str, content: str, metadata: Dict | None = None) -> Dict:
    normalized = _normalize_text(content)
    return {
        'chunk_index': index,
        'source_page': source_page,
        'heading': _normalize_text(heading),
        'content': normalized,
        'keyword_summary': _summarize_keywords(normalized),
        'metadata': metadata or {},
    }


def parse_pdf(file_path: str) -> Dict:
    try:
        from pypdf import PdfReader
    except Exception as exc:
        raise MaterialParseError('缺少 pypdf 依赖，暂时无法解析 PDF') from exc

    try:
        reader = PdfReader(file_path)
    except Exception as exc:
        raise MaterialParseError(f'PDF 打开失败: {exc}') from exc

    chunks: List[Dict] = []
    full_text_parts: List[str] = []
    for page_index, page in enumerate(reader.pages, start=1):
        try:
            text = _normalize_text(page.extract_text() or '')
        except Exception:
            text = ''
        if not text:
            continue
        full_text_parts.append(text)
        chunks.append(_build_chunk(len(chunks), str(page_index), f'第 {page_index} 页', text, {'page_number': page_index}))

    if not chunks:
        raise MaterialParseError('PDF 未提取到可用文本')

    return {
        'page_count': len(reader.pages),
        'full_text': '\n\n'.join(full_text_parts),
        'chunks': chunks,
    }


def parse_pptx(file_path: str) -> Dict:
    try:
        from pptx import Presentation
    except Exception as exc:
        raise MaterialParseError('缺少 python-pptx 依赖，暂时无法解析 PPTX') from exc

    if file_path.lower().endswith('.ppt'):
        raise MaterialParseError('当前仅支持解析 .pptx 文件，旧版 .ppt 请先另存为 .pptx')

    try:
        presentation = Presentation(file_path)
    except Exception as exc:
        raise MaterialParseError(f'PPTX 打开失败: {exc}') from exc

    chunks: List[Dict] = []
    full_text_parts: List[str] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        texts = []
        for shape in slide.shapes:
            try:
                if hasattr(shape, 'text') and shape.text:
                    texts.append(shape.text)
            except Exception:
                continue
        slide_text = _normalize_text('\n'.join(texts))
        if not slide_text:
            continue
        lines = [line.strip() for line in slide_text.splitlines() if line.strip()]
        heading = lines[0] if lines else f'第 {slide_index} 页'
        full_text_parts.append(slide_text)
        chunks.append(_build_chunk(len(chunks), str(slide_index), heading, slide_text, {'slide_number': slide_index}))

    if not chunks:
        raise MaterialParseError('PPTX 未提取到可用文本')

    return {
        'page_count': len(presentation.slides),
        'full_text': '\n\n'.join(full_text_parts),
        'chunks': chunks,
    }


def parse_material_file(file_path: str, material_type: str = '', file_name: str = '') -> Dict:
    extension = os.path.splitext(file_name or file_path)[1].lower()
    kind = (material_type or '').lower()

    if kind == 'pdf' or extension == '.pdf':
        return parse_pdf(file_path)
    if kind == 'ppt' or extension in {'.pptx', '.ppt'}:
        return parse_pptx(file_path)

    raise MaterialParseError('当前仅实现 PDF 和 PPTX 的解析，其他格式后续接入')
