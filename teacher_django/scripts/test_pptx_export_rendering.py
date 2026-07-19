#!/usr/bin/env python
"""
直接测试 pptx_exporter.export_outline_to_pptx 的渲染效果。

手写一份包含 quiz_check（双题）、two_column 中的 code 块、animation_embed
（concept_name/usage_note）、comparison、process_flow 等布局的 structured_slides，
调用导出函数生成 .pptx 文件，再用 python-pptx 读回所有文本，断言关键内容均出现。

同时测试 structured_slides 为空时的兜底路径（normalize_slide_deck(None, ...)）。
"""
import os
import sys
import json

try:
    sys.stdout.reconfigure(encoding='utf-8')
except Exception:
    pass

BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
if BASE_DIR not in sys.path:
    sys.path.insert(0, BASE_DIR)

os.environ.setdefault('DJANGO_SETTINGS_MODULE', 'teacher_django.settings')
import django
django.setup()

from django.contrib.auth import get_user_model
from curriculum_app.models import CourseOutline
from curriculum_app.utils.pptx_exporter import export_outline_to_pptx

User = get_user_model()


def collect_all_text(pptx_path):
    from pptx import Presentation
    prs = Presentation(pptx_path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text:
                            texts.append(run.text)
    return texts, len(prs.slides._sldIdLst)


def build_sample_outline():
    structured_slides = [
        {
            'layout': 'cover',
            'theme': 'academic_light',
            'title': '梯度下降入门',
            'teaching_goal': '理解梯度下降的基本原理',
            'bullets': ['梯度下降是优化算法', '用于最小化损失函数', '是深度学习训练的基础'],
            'speaker_notes': '欢迎大家学习梯度下降！',
        },
        {
            'layout': 'two_column',
            'theme': 'tech_blue',
            'title': '梯度下降代码实现',
            'teaching_goal': '掌握 Python 实现方式',
            'bullets': ['核心是沿负梯度方向更新参数'],
            'visual_blocks': [
                {
                    'kind': 'code',
                    'label': '梯度下降实现',
                    'language': 'python',
                    'code': "\n".join([
                        "def gradient_descent(x0, lr=0.1, steps=50):",
                        "    x = x0",
                        "    for i in range(steps):",
                        "        grad = 2 * x  # f(x)=x^2 的导数",
                        "        x = x - lr * grad",
                        "    return x",
                        "",
                        "print(gradient_descent(10.0))",
                    ]),
                },
                {
                    'kind': 'bullet_card',
                    'label': '关键参数',
                    'items': ['学习率 lr 控制步长', '迭代次数 steps 控制训练轮数'],
                },
            ],
            'speaker_notes': '这是一个非常长的讲稿用于测试截断效果。' * 10,
        },
        {
            'layout': 'animation_embed',
            'theme': 'chalkboard_dark',
            'title': '梯度下降动画演示',
            'teaching_goal': '直观理解收敛过程',
            'bullets': [],
            'visual_blocks': [
                {
                    'kind': 'animation',
                    'label': '梯度下降收敛过程',
                    'concept_name': '梯度下降收敛过程',
                    'usage_note': '观察小球如何沿曲面滚向最低点，体会学习率对收敛速度的影响。',
                    'animation_code': '<!doctype html><html><body>demo</body></html>',
                },
            ],
            'student_interaction': '请预测：学习率过大会发生什么？',
            'speaker_notes': '播放动画并引导学生观察。',
        },
        {
            'layout': 'quiz_check',
            'theme': 'academic_light',
            'title': '随堂检测：梯度下降是否正确这是一个非常非常长的标题用于测试标题溢出场景下的自动截断与缩放效果',
            'teaching_goal': '检验学习效果',
            'bullets': [],
            'visual_blocks': [
                {
                    'kind': 'question',
                    'label': '问题1',
                    'question_text': '梯度下降沿什么方向更新参数？',
                    'question_type': 'choice',
                    'choices': [
                        {'label': 'A. 梯度方向', 'value': 'A'},
                        {'label': 'B. 负梯度方向', 'value': 'B'},
                        {'label': 'C. 随机方向', 'value': 'C'},
                        {'label': 'D. 零方向', 'value': 'D'},
                    ],
                    'correct_answer': 'B',
                    'explanation': '负梯度方向是损失函数下降最快的方向，因此参数沿该方向更新。',
                },
                {
                    'kind': 'question',
                    'label': '问题2',
                    'question_text': '学习率过大可能导致什么问题？',
                    'question_type': 'choice',
                    'choices': [
                        {'label': 'A. 收敛太慢', 'value': 'A'},
                        {'label': 'B. 可能发散', 'value': 'B'},
                        {'label': 'C. 没有问题', 'value': 'C'},
                    ],
                    'correct_answer': 'B',
                    'explanation': '学习率过大会导致参数更新步长过大，可能跳过最优解甚至发散。',
                },
            ],
            'speaker_notes': '让学生先独立作答，再公布答案。',
        },
        {
            'layout': 'comparison',
            'theme': 'tech_blue',
            'title': '批量梯度下降 vs 随机梯度下降',
            'teaching_goal': '理解两种方法的差异',
            'bullets': [],
            'visual_blocks': [
                {'kind': 'compare_column', 'label': '批量梯度下降', 'items': ['每次用全部样本计算梯度', '收敛稳定但速度慢']},
                {'kind': 'compare_column', 'label': '随机梯度下降', 'items': ['每次用单个样本计算梯度', '速度快但波动大']},
            ],
            'speaker_notes': '对比两种方法的优缺点。',
        },
        {
            'layout': 'summary',
            'theme': 'academic_light',
            'title': '本讲小结',
            'teaching_goal': '回顾核心知识点',
            'bullets': ['梯度下降原理', '学习率的作用', '批量与随机梯度下降的差异'],
            'speaker_notes': '总结今天的内容。',
        },
    ]

    outline_data = {
        'title': '梯度下降专题课',
        'resources': {
            'ppt': {
                'structured_slides': structured_slides,
            }
        },
    }
    return outline_data


def main():
    user, _ = User.objects.get_or_create(username='pptx_render_test', defaults={'email': 'pptx_render_test@example.com'})

    outline_data = build_sample_outline()
    outline = CourseOutline.objects.create(
        user=user,
        title='梯度下降专题课',
        outline_data=json.dumps(outline_data, ensure_ascii=False),
        status='completed',
        progress=100,
    )
    print('Created outline id', outline.id)

    file_path, filename = export_outline_to_pptx(outline)
    print('Exported file:', file_path)
    assert os.path.exists(file_path), 'PPTX 文件未生成'
    assert os.path.getsize(file_path) > 0, 'PPTX 文件为空'

    texts, slide_count = collect_all_text(file_path)
    full_text = '\n'.join(texts)
    print(f'共 {slide_count} 页幻灯片，{len(texts)} 个文本片段')

    # 期望页数 = 封面(add_deck_cover) + 6 个 structured_slides
    assert slide_count == 7, f'预期 7 页（1封面 + 6内容页），实际 {slide_count} 页'

    checks = {
        '代码块标题': '梯度下降实现',
        '代码内容': 'def gradient_descent',
        '代码内容2': 'grad = 2 * x',
        '代码截断标记或正常结尾': 'print(gradient_descent(10.0))',
        '动画概念名': '梯度下降收敛过程',
        '动画使用说明': '观察小球如何沿曲面滚向最低点',
        '动画跳转提示': '请前往网页课件查看交互动画演示',
        '题目1': '梯度下降沿什么方向更新参数？',
        '题目1正确选项': '负梯度方向',
        '题目1解析': '负梯度方向是损失函数下降最快的方向',
        '题目2': '学习率过大可能导致什么问题？',
        '题目2正确选项': '可能发散',
        '对比列1': '批量梯度下降',
        '对比列2': '随机梯度下降',
        '小结要点': '学习率的作用',
    }
    failed = []
    for name, expect in checks.items():
        if expect not in full_text:
            failed.append((name, expect))

    if failed:
        print('\n[FAIL] 以下内容未在导出的 PPTX 中找到：')
        for name, expect in failed:
            print(f'  - {name}: {expect!r}')
        print('\n--- 全部文本片段 ---')
        for t in texts:
            print(repr(t))
        raise SystemExit(1)

    print('[OK] 所有结构化内容均已正确渲染到 PPTX 中')

    # 标题溢出截断检查：超长标题应被截断为 40 字符 + ...
    long_title_found = any(t.startswith('随堂检测：梯度下降') and t.endswith('...') for t in texts)
    assert long_title_found, '超长标题未按预期截断'
    print('[OK] 超长标题已按预期截断')

    # 讲稿截断检查
    long_notes_found = any(t.startswith('讲稿：') and t.endswith('...') and len(t) < 130 for t in texts)
    assert long_notes_found, '超长讲稿未按预期截断'
    print('[OK] 超长讲稿已按预期截断')

    print('\n=== 测试 structured_slides 缺失时的兜底路径 ===')
    fallback_outline = CourseOutline.objects.create(
        user=user,
        title='兜底测试课程：函数与映射',
        outline_data=json.dumps({'title': '兜底测试课程：函数与映射'}, ensure_ascii=False),
        status='completed',
        progress=100,
    )
    fb_path, fb_name = export_outline_to_pptx(fallback_outline)
    assert os.path.exists(fb_path) and os.path.getsize(fb_path) > 0
    fb_texts, fb_slide_count = collect_all_text(fb_path)
    assert fb_slide_count >= 4, f'兜底路径应至少生成 1 封面 + 3 内容页，实际 {fb_slide_count} 页'
    print(f'[OK] 兜底路径生成 {fb_slide_count} 页 PPTX')

    print('\n[SUCCESS] PPTX 渲染测试全部通过！')


if __name__ == '__main__':
    main()
